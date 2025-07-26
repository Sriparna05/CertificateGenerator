import os
import shutil
import uuid
import logging
import platform
from datetime import datetime

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    Image = None

from jinja2 import Environment, FileSystemLoader
from app.celery_worker import celery_app

# --- START: WEASYPRINT & PDF2IMAGE DEPENDENCY LOADER ---

WEASYPRINT_INSTALLED = False
PDF2IMAGE_INSTALLED = False

def load_dependencies():
    """Manually loads WeasyPrint's GTK dependencies and checks for other libraries."""
    global WEASYPRINT_INSTALLED, PDF2IMAGE_INSTALLED

    if platform.system() == "Windows":
        try:
            # NOTE: Update this path if your GTK3 installation is different.
            gtk_path = r"C:\Program Files\GTK3-Runtime Win64\bin"
            if os.path.exists(gtk_path):
                os.add_dll_directory(gtk_path)
                logging.debug(f"Successfully added GTK3 DLL directory to search path: {gtk_path}")
            else:
                logging.warning(f"GTK3 installation not found at default path: {gtk_path}. PDF generation will likely fail.")
        except Exception as e:
            logging.error(f"Could not add GTK3 DLL directory. Ensure you are on Python 3.8+. Error: {e}")

    try:
        import weasyprint
        WEASYPRINT_INSTALLED = True
        logging.info("WeasyPrint is available.")
    except ImportError:
        logging.error("WeasyPrint not found. To generate PDFs, please run 'pip install weasyprint'.")
        WEASYPRINT_INSTALLED = False
        
    try:
        import pdf2image
        PDF2IMAGE_INSTALLED = True
        logging.info("pdf2image is available.")
    except ImportError:
        logging.warning("pdf2image not found. To generate Images from HTML, please run 'pip install pdf2image' and install Poppler.")
        PDF2IMAGE_INSTALLED = False

load_dependencies()

# --- END: DEPENDENCY LOADER ---


@celery_app.task
def generate_certificate_async(data):
    """Celery task to generate certificates asynchronously."""
    try:
        return generate_certificate(
            data.get("template_id"),
            data.get("output_format", "pdf"),
            data.get("recipients", []),
            data.get("ai_options", {}),
        )
    except Exception as e:
        return {"status": "error", "message": f"Async generation failed: {str(e)}"}


def list_templates(base_dir="certificate_templates"):
    """List available certificate templates by type."""
    templates = {"pptx": [], "images": [], "html": []}
    for type_folder in templates.keys():
        folder_path = os.path.join(base_dir, type_folder)
        if os.path.isdir(folder_path):
            templates[type_folder] = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
    return templates


def html_to_pdf(html_path, pdf_output_path):
    """Convert HTML certificate to PDF with a file existence check."""
    if not WEASYPRINT_INSTALLED:
        raise ImportError("Cannot generate PDF. WeasyPrint is not properly installed or its system dependencies (GTK3) are missing.")

    from weasyprint import HTML, CSS
    custom_css = CSS(string="@page { size: A4 landscape; margin: 0; }")
    HTML(filename=html_path).write_pdf(pdf_output_path, stylesheets=[custom_css])
    
    # --- ADDED ROBUSTNESS CHECK ---
    if not os.path.exists(pdf_output_path):
        raise FileNotFoundError(f"WeasyPrint ran but failed to create the PDF file at: {pdf_output_path}")

    return pdf_output_path


def html_to_image(html_path, image_output_path, format="PNG"):
    """Convert HTML certificate to an image."""
    if not PDF2IMAGE_INSTALLED:
        raise ImportError("Cannot generate images. pdf2image or its dependency (Poppler) is not installed.")

    temp_pdf_path = image_output_path.replace(f".{format.lower()}", "_temp.pdf")
    html_to_pdf(html_path, temp_pdf_path)

    from pdf2image import convert_from_path
    images = convert_from_path(temp_pdf_path, dpi=300, first_page=1, last_page=1)
    
    if not images:
        os.remove(temp_pdf_path)
        raise Exception("PDF to image conversion resulted in no images.")
    
    images[0].save(image_output_path, format=format, quality=95)
    os.remove(temp_pdf_path)
    return image_output_path


def generate_certificate_from_html(template_dir, template_name, output_path, context):
    """Render an HTML template with context and save to output_path."""
    env = Environment(loader=FileSystemLoader(template_dir))
    template = env.get_template(template_name)
    rendered_html = template.render(context)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(rendered_html)


def generate_certificate_from_pptx(template_path, output_path, replacements):
    """Generate a certificate from a PPTX template by replacing placeholders."""
    if not Presentation: raise ImportError("python-pptx library is not installed.")
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                # This complex replacement logic attempts to preserve formatting.
                # A simpler version would be: p.text = p.text.replace(key, str(value))
                full_text = "".join(run.text for run in p.runs)
                if any(key in full_text for key in replacements):
                    original_runs = list(p.runs)
                    p.clear() # Clear existing runs
                    
                    temp_text = full_text
                    for key, value in replacements.items():
                        temp_text = temp_text.replace(key, str(value))
                    
                    new_run = p.add_run()
                    new_run.text = temp_text
                    if original_runs:
                        new_run.font.bold = original_runs[0].font.bold
                        new_run.font.italic = original_runs[0].font.italic
                        new_run.font.name = original_runs[0].font.name
                        new_run.font.size = original_runs[0].font.size
                        new_run.font.color.rgb = original_runs[0].font.color.rgb
    prs.save(output_path)


def generate_certificate(template_id, output_format, recipients, ai_options=None):
    """Generate certificates for multiple recipients."""
    results = []
    template_base_dir = "certificate_templates"
    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_dir = os.path.join(project_root, "generated_certificates")
    os.makedirs(output_dir, exist_ok=True)

    for recipient in recipients:
        certificate_id = str(uuid.uuid4())
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        recipient_name = recipient.get("name", "Unknown")
        final_path = None
        error_message = None

        try:
            logging.debug(f"Processing {recipient_name} for template {template_id}")

            if template_id.endswith(".html"):
                template_dir = os.path.join(template_base_dir, "html")
                html_path = os.path.join(output_dir, f"cert_{certificate_id}_{timestamp}.html")
                context = {**recipient, "certificate_id": certificate_id}
                
                generate_certificate_from_html(template_dir, template_id, html_path, context)

                if output_format.lower() == "pdf":
                    pdf_path = html_path.replace(".html", ".pdf")
                    final_path = html_to_pdf(html_path, pdf_path)
                elif output_format.lower() in ["png", "jpeg"]:
                    image_path = html_path.replace(".html", f".{output_format.lower()}")
                    final_path = html_to_image(html_path, image_path, output_format.upper())
                else:
                    final_path = html_path

            elif template_id.endswith(".pptx"):
                if output_format.lower() != "pptx":
                    raise ValueError("PPTX templates can only be generated in PPTX format.")
                
                template_path = os.path.join(template_base_dir, "pptx", template_id)
                pptx_path = os.path.join(output_dir, f"cert_{certificate_id}_{timestamp}.pptx")
                
                replacements = {"{{%s}}" % key: str(value) for key, value in recipient.items()}
                generate_certificate_from_pptx(template_path, pptx_path, replacements)
                final_path = pptx_path

            else:
                raise ValueError(f"Unsupported template format: {template_id}")

        except Exception as e:
            logging.error(f"Failed to generate certificate for {recipient_name}: {e}", exc_info=True)
            error_message = str(e)

        relative_path = os.path.relpath(final_path, project_root).replace("\\", "/") if final_path and os.path.exists(final_path) else None
        
        results.append({
            "recipient": recipient_name,
            "certificate_id": certificate_id,
            "file_path": relative_path,
            "status": "success" if relative_path else "error",
            "error": error_message,
        })

    successful_count = sum(1 for r in results if r["status"] == "success")
    return {
        "status": "completed",
        "total_recipients": len(recipients),
        "successful": successful_count,
        "failed": len(recipients) - successful_count,
        "results": results,
    }